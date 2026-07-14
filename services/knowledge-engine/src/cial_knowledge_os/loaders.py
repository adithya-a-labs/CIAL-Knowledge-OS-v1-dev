"""Local document ingestion with a registry-driven enterprise corpus."""

from __future__ import annotations

import csv
import json
import logging
import mimetypes
import time
from collections import Counter
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from langchain_core.documents import Document

from .config import KnowledgeOSConfig
from .file_formats import (
    SupportStatus,
    get_file_format_info,
    is_supported_file,
    inspect_ingestion_candidate,
    list_supported_formats,
    scan_file_format_readiness,
    validate_ingestion_file,
)
from .ocr import create_ocr_engine

logger = logging.getLogger(__name__)

SUPPORTED_DOCUMENT_EXTENSIONS = frozenset(
    f".{extension}"
    for definition in list_supported_formats()
    for extension in definition["extensions"]
)
IMPLEMENTED_KNOWLEDGE_EXTENSIONS = SUPPORTED_DOCUMENT_EXTENSIONS
TEXT_LIKE_EXTENSIONS = frozenset(
    {".txt", ".md", ".markdown", ".html", ".htm", ".json", ".xml", ".yaml", ".yml", ".csv"}
)
DOCLING_EXTENSIONS = frozenset(
    {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"}
)
OCR_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".tiff", ".tif"})

SAMPLE_AIRPORT_DOCUMENTS = {
    "terminal_operations_sop.txt": """CIAL Terminal Operations SOP
Version: 1.0
Owner: Airport Operations Control Centre

Passenger queue monitoring must be reviewed every 15 minutes during peak periods.
Duty managers must coordinate with airline supervisors when a queue exceeds the
marked holding area. Escalators, elevators, baggage belts, and passenger boarding
bridges must be visually checked at the start of each shift.
""",
    "runway_maintenance_sop.txt": """CIAL Runway Maintenance SOP
Version: 1.0
Owner: Airside Operations

Routine runway surface inspection must be performed every 6 hours and after heavy
rain, bird-strike reports, foreign-object-debris alerts, or pilot braking-action
complaints. Maintenance teams must obtain ATC clearance before entering the runway
strip and maintain continuous radio contact.
""",
    "electrical_maintenance_manual.txt": """CIAL Electrical Maintenance Manual
Version: 1.0
Owner: Electrical Engineering

Work on energized panels requires voltage-rated insulated gloves, an arc-rated face
shield, flame-resistant clothing, dielectric safety shoes, insulated tools, and
lockout-tagout verification wherever isolation is possible. Emergency shutdown
requires operating the designated stop, notifying electrical control, applying
lockout-tagout, and recording the event.
""",
}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = " ".join(data.split())
        if text:
            self.parts.append(text)

    def text(self) -> str:
        return "\n".join(self.parts)


def create_sample_airport_documents(config: KnowledgeOSConfig) -> list[Path]:
    """Explicitly create non-sensitive fixtures without overwriting edits."""

    config.sample_data_dir.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    for file_name, content in SAMPLE_AIRPORT_DOCUMENTS.items():
        path = config.sample_data_dir / file_name
        if not path.exists():
            path.write_text(content.strip() + "\n", encoding="utf-8")
        paths.append(path)
    return paths


def _base_metadata(
    path: Path,
    loader_type: str,
    page_number: int | None = None,
    *,
    corpus_root: Path | None = None,
    repository_id: str | None = None,
) -> dict[str, Any]:
    resolved = path.resolve()
    relative_path: Path | None = None
    if corpus_root is not None:
        try:
            relative_path = resolved.relative_to(corpus_root.resolve())
        except ValueError:
            relative_path = None
    folder_parts = relative_path.parts[:-1] if relative_path is not None else ()
    format_info = get_file_format_info(path.name)
    metadata: dict[str, Any] = {
        "source": str(resolved),
        "file_name": path.name,
        "source_filename": path.name,
        "absolute_path": str(resolved),
        "relative_path": relative_path.as_posix() if relative_path else path.name,
        "repository_id": repository_id,
        "category": folder_parts[0] if folder_parts else None,
        "collection": folder_parts[1] if len(folder_parts) > 1 else None,
        "loader_type": loader_type,
        "document_type": format_info["extension"],
        "file_type": path.suffix.lstrip(".").casefold(),
        "mime_type": mimetypes.guess_type(path.name)[0],
        "file_format_category": format_info["category"],
        "file_format_label": format_info["format_label"],
        "file_support_status": format_info["support_status"],
        "requires_ocr": bool(format_info["requires_ocr"]),
        "access_level": "internal",
    }
    if page_number is not None:
        metadata["page_number"] = page_number
        metadata["page_index"] = page_number - 1
        metadata["citation_metadata_version"] = 2
    return metadata


def _read_text_lossy(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


def _load_text_like_document(
    path: Path,
    corpus_root: Path | None = None,
    *,
    repository_id: str | None = None,
) -> list[Document]:
    suffix = path.suffix.casefold()
    raw = _read_text_lossy(path)
    text = raw
    loader_type = "text"
    if suffix in {".html", ".htm"}:
        parser = _HTMLTextExtractor()
        parser.feed(raw)
        text = parser.text()
        loader_type = "html"
    elif suffix == ".json":
        try:
            text = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            text = raw
        loader_type = "json"
    elif suffix == ".csv":
        try:
            rows = list(csv.reader(raw.splitlines()))
            text = "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
        except csv.Error:
            text = raw
        loader_type = "csv"
    elif suffix in {".yaml", ".yml"}:
        loader_type = "yaml"
    elif suffix == ".xml":
        loader_type = "xml"
    elif suffix in {".md", ".markdown"}:
        loader_type = "markdown"
    text = text.strip()
    if not text:
        return []
    return [
        Document(
            page_content=text,
            metadata=_base_metadata(
                path,
                loader_type,
                corpus_root=corpus_root,
                repository_id=repository_id,
            ),
        )
    ]


def load_text_documents(config: KnowledgeOSConfig) -> list[Document]:
    """Load UTF-8 text fixtures from sample and raw local data directories."""

    documents: list[Document] = []
    seen: set[Path] = set()
    for directory in (config.sample_data_dir, config.raw_data_dir):
        if not directory.exists():
            continue
        for path in sorted(directory.rglob("*.txt")):
            resolved = path.resolve()
            if resolved in seen:
                continue
            seen.add(resolved)
            documents.extend(
                _load_text_like_document(
                    path,
                    corpus_root=config.knowledge_root,
                    repository_id=config.repository_id,
                )
            )
    return documents


def _load_with_docling(path: Path, corpus_root: Path, *, repository_id: str | None = None) -> list[Document]:
    from docling.document_converter import DocumentConverter

    result = DocumentConverter().convert(str(path))
    text = result.document.export_to_markdown().strip()
    if not text:
        return []
    return [
        Document(
            page_content=text,
            metadata=_base_metadata(path, "docling", corpus_root=corpus_root, repository_id=repository_id),
        )
    ]


def _load_pdf_with_docling(path: Path, corpus_root: Path, *, repository_id: str | None = None) -> list[Document]:
    return _load_with_docling(path, corpus_root, repository_id=repository_id)


def _spreadsheet_row_text(values: list[object]) -> str:
    cells = [str(value).strip() for value in values if value not in {None, ""}]
    return " | ".join(cell for cell in cells if cell)


def _load_xlsx_document(path: Path, corpus_root: Path, *, repository_id: str | None = None) -> list[Document]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ImportError("openpyxl is required for XLSX ingestion.") from exc

    documents: list[Document] = []
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
            lines: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                row_text = _spreadsheet_row_text(list(row))
                if row_text:
                    lines.append(row_text)
            text = "\n".join(lines).strip()
            if not text:
                continue
            metadata = _base_metadata(path, "xlsx", corpus_root=corpus_root, repository_id=repository_id)
            metadata["sheet_name"] = worksheet.title
            metadata["sheet_index"] = sheet_index
            metadata["anchor"] = f"sheet:{worksheet.title}"
            documents.append(Document(page_content=text, metadata=metadata))
    finally:
        workbook.close()
    return documents


def _load_xls_document(path: Path, corpus_root: Path, *, repository_id: str | None = None) -> list[Document]:
    try:
        import xlrd
    except ImportError as exc:
        raise ImportError("xlrd is required for XLS ingestion.") from exc

    documents: list[Document] = []
    workbook = xlrd.open_workbook(path.as_posix(), on_demand=True)
    try:
        for sheet_index, sheet_name in enumerate(workbook.sheet_names(), start=1):
            sheet = workbook.sheet_by_index(sheet_index - 1)
            lines: list[str] = []
            for row_index in range(sheet.nrows):
                row_text = _spreadsheet_row_text(
                    [sheet.cell_value(row_index, col_index) for col_index in range(sheet.ncols)]
                )
                if row_text:
                    lines.append(row_text)
            text = "\n".join(lines).strip()
            if not text:
                continue
            metadata = _base_metadata(path, "xls", corpus_root=corpus_root, repository_id=repository_id)
            metadata["sheet_name"] = sheet_name
            metadata["sheet_index"] = sheet_index
            metadata["anchor"] = f"sheet:{sheet_name}"
            documents.append(Document(page_content=text, metadata=metadata))
    finally:
        workbook.release_resources()
    return documents


def _load_pptx_document(path: Path, corpus_root: Path, *, repository_id: str | None = None) -> list[Document]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("python-pptx is required for PPTX ingestion.") from exc

    documents: list[Document] = []
    presentation = Presentation(path)
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
        slide_text = "\n".join(parts).strip()
        if not slide_text:
            continue
        metadata = _base_metadata(path, "pptx", corpus_root=corpus_root, repository_id=repository_id)
        metadata["slide_number"] = slide_number
        metadata["anchor"] = f"slide:{slide_number}"
        documents.append(Document(page_content=slide_text, metadata=metadata))
    return documents


def _load_pdf_with_pymupdf(path: Path, corpus_root: Path, *, repository_id: str | None = None) -> list[Document]:
    import fitz

    documents: list[Document] = []
    with fitz.open(path) as pdf:
        for index, page in enumerate(pdf, start=1):
            text = page.get_text("text").strip()
            if text:
                documents.append(
                    Document(
                        page_content=text,
                        metadata=_base_metadata(
                            path,
                            "pymupdf",
                            index,
                            corpus_root=corpus_root,
                            repository_id=repository_id,
                        ),
                    )
                )
    return documents


def _load_ocr_document(path: Path, config: KnowledgeOSConfig, corpus_root: Path) -> list[Document]:
    inspection = inspect_ingestion_candidate(path, corpus_root=corpus_root, ocr_engine=config.ocr_engine)
    if not config.ocr_enabled:
        logger.warning(
            "ocr_skipped_disabled",
            extra={
                "event": "ocr",
                "document_filename": inspection["filename"],
                "relative_path": inspection["relative_path"],
                "extension": inspection["extension"],
                "detected_mime_type": inspection["detected_mime_type"],
                "file_size_bytes": inspection["file_size_bytes"],
                "ocr_engine": config.ocr_engine,
                "ocr_engine_available": False,
                "ocr_fallback_available": False,
                "skip_reason": "OCR_disabled_by_configuration",
            },
        )
        return []
    engine = create_ocr_engine(config)
    preflight = getattr(engine, "preflight", lambda **_: {"status": "unknown"})(enabled=True)
    result = engine.extract(path)
    metadata = _base_metadata(path, "ocr", corpus_root=corpus_root, repository_id=config.repository_id)
    metadata.update(result.metadata)
    metadata["ocr_status"] = result.status
    if result.error:
        metadata["ocr_error"] = result.error
    if result.status != "OCR_SUCCESS" or not result.text.strip():
        logger.warning(
            "ocr_file_skipped",
            extra={
                "event": "ocr",
                "document_filename": inspection["filename"],
                "relative_path": inspection["relative_path"],
                "extension": inspection["extension"],
                "detected_mime_type": inspection["detected_mime_type"],
                "file_size_bytes": inspection["file_size_bytes"],
                "ocr_engine": config.ocr_engine,
                "ocr_engine_available": preflight.get("status") == "ready",
                "ocr_fallback_available": False,
                "skip_reason": "image_decode_failed" if result.status == "OCR_FAILED" else "already_contains_extractable_text",
                "error_type": type(result.error).__name__ if result.error else None,
                "ocr_status": result.status,
                "error": result.error[:300] if result.error else None,
            },
        )
        return []
    return [Document(page_content=result.text, metadata=metadata)]


def _supported_documents(root: Path) -> list[Path]:
    if not root.exists():
        return []
    return sorted(path for path in root.rglob("*") if path.is_file() and is_supported_file(path.name))


def resolve_corpus_root(config: KnowledgeOSConfig) -> Path:
    """Return the canonical configured corpus root."""

    return config.knowledge_root


def discover_knowledge_documents(
    config: KnowledgeOSConfig,
) -> tuple[Path, list[Path]]:
    """Recursively find processable files below the active configured corpus."""

    corpus_root = resolve_corpus_root(config)
    if not corpus_root.exists():
        return corpus_root, []
    paths: list[Path] = []
    for path in sorted(item for item in corpus_root.rglob("*") if item.is_file()):
        inspection = inspect_ingestion_candidate(path, corpus_root=corpus_root, ocr_engine=config.ocr_engine)
        validation = inspection["validation"]
        if inspection["eligible"] and validation["valid_for_ingestion"]:
            paths.append(path)
            continue
        logger.warning(
            "document_type_not_ingested",
            extra={
                "event": "document_discovery",
                "document_filename": inspection["filename"],
                "relative_path": inspection["relative_path"],
                "extension": inspection["extension"],
                "detected_mime_type": inspection["detected_mime_type"],
                "configured_supported_formats": sorted(SUPPORTED_DOCUMENT_EXTENSIONS),
                "skip_reason": inspection["skip_reason"],
                "loader_selected": inspection["loader_selected"],
                "support_status": validation["support_status"],
            },
        )
    return corpus_root, paths


def _load_supported_path(path: Path, config: KnowledgeOSConfig, corpus_root: Path) -> list[Document]:
    inspection = inspect_ingestion_candidate(path, corpus_root=corpus_root, ocr_engine=config.ocr_engine)
    validation = inspection["validation"]
    if not inspection["eligible"] or not validation["valid_for_ingestion"]:
        logger.warning(
            "document_type_not_ingested",
            extra={
                "event": "document_loading",
                "document_filename": inspection["filename"],
                "relative_path": inspection["relative_path"],
                "extension": inspection["extension"],
                "detected_mime_type": inspection["detected_mime_type"],
                "configured_supported_formats": sorted(SUPPORTED_DOCUMENT_EXTENSIONS),
                "skip_reason": inspection["skip_reason"],
                "loader_selected": inspection["loader_selected"],
                "support_status": validation["support_status"],
            },
        )
        return []
    suffix = path.suffix.casefold()
    if suffix in OCR_EXTENSIONS:
        return _load_ocr_document(path, config, corpus_root)
    if suffix in TEXT_LIKE_EXTENSIONS:
        return _load_text_like_document(path, corpus_root, repository_id=config.repository_id)
    if suffix == ".pdf":
        return _load_pdf_path(path, corpus_root, repository_id=config.repository_id)
    if suffix == ".xlsx":
        try:
            documents = _load_xlsx_document(path, corpus_root, repository_id=config.repository_id)
            if documents:
                return documents
        except Exception as exc:
            logger.warning(
                "xlsx_document_load_failed",
                extra={
                    "event": "document_loading",
                    "source": str(path.resolve()),
                    "error": str(exc),
                },
            )
    if suffix == ".xls":
        try:
            documents = _load_xls_document(path, corpus_root, repository_id=config.repository_id)
            if documents:
                return documents
        except Exception as exc:
            logger.warning(
                "xls_document_load_failed",
                extra={
                    "event": "document_loading",
                    "source": str(path.resolve()),
                    "error": str(exc),
                },
            )
    if suffix == ".pptx":
        try:
            documents = _load_pptx_document(path, corpus_root, repository_id=config.repository_id)
            if documents:
                return documents
        except Exception as exc:
            logger.warning(
                "pptx_document_load_failed",
                extra={
                    "event": "document_loading",
                    "source": str(path.resolve()),
                    "error": str(exc),
                },
            )
    if suffix in DOCLING_EXTENSIONS:
        try:
            return _load_with_docling(path, corpus_root, repository_id=config.repository_id)
        except ImportError as exc:
            logger.warning(
                "docling_document_loader_unavailable",
                extra={
                    "event": "document_loading",
                    "source": str(path.resolve()),
                    "error": str(exc),
                },
            )
            return []
        except Exception as exc:
            logger.warning(
                "docling_document_load_failed",
                extra={
                    "event": "document_loading",
                    "source": str(path.resolve()),
                    "error": str(exc),
                },
            )
            return []
    return []


def _load_pdf_path(path: Path, corpus_root: Path, *, repository_id: str | None = None) -> list[Document]:
    try:
        import docling  # noqa: F401
    except ImportError:
        docling_available = False
    else:
        docling_available = True

    try:
        import fitz  # noqa: F401
    except ImportError:
        pymupdf_available = False
    else:
        pymupdf_available = True

    if not docling_available and not pymupdf_available:
        raise ImportError(
            "PDF ingestion requires a local loader. Install 'docling' (preferred) "
            "or 'PyMuPDF'; no cloud OCR fallback is used."
        )

    # PDF citations must be page-addressable.  Docling's current Markdown
    # export is document-wide and does not retain page provenance, whereas
    # PyMuPDF emits one Document per physical PDF page. Prefer the latter when
    # available; Docling remains a local fallback for environments without it.
    if pymupdf_available:
        try:
            return (
                _load_pdf_with_pymupdf(path, corpus_root, repository_id=repository_id)
                if repository_id is not None
                else _load_pdf_with_pymupdf(path, corpus_root)
            )
        except Exception as exc:
            if not docling_available:
                raise RuntimeError(
                    f"Could not read PDF '{path}'. The file may be corrupted, "
                    f"encrypted, or unsupported by PyMuPDF. Original error: {exc}"
                ) from exc
            logger.warning(
                "pymupdf_pdf_fallback",
                extra={"event": "document_loading", "source": str(path), "error": str(exc)},
            )

    if docling_available:
        try:
            docling_documents = (
                _load_pdf_with_docling(path, corpus_root, repository_id=repository_id)
                if repository_id is not None
                else _load_pdf_with_docling(path, corpus_root)
            )
            if docling_documents:
                return docling_documents
        except Exception as exc:
            if not pymupdf_available:
                raise RuntimeError(f"Docling could not load {path.name}: {exc}") from exc
            logger.warning(
                "docling_pdf_fallback",
                extra={
                    "event": "document_loading",
                    "source": str(path),
                    "error": str(exc),
                },
            )
        else:
            if not pymupdf_available:
                raise RuntimeError(
                    f"Docling extracted no text from {path.name}, and PyMuPDF "
                    "is not installed as a local fallback."
                )
    raise RuntimeError(
        f"No text could be extracted from PDF '{path}'. Docling output does not "
        "include page provenance, so this PDF cannot be indexed for page-addressable citations."
    )


def load_pdf_documents(config: KnowledgeOSConfig) -> list[Document]:
    """Recursively load all processable enterprise corpus documents."""

    corpus_root, paths = discover_knowledge_documents(config)
    return load_pdf_paths(paths, corpus_root=corpus_root, config=config)


def load_pdf_paths(
    pdf_paths: list[Path],
    *,
    corpus_root: Path,
    config: KnowledgeOSConfig | None = None,
) -> list[Document]:
    """Load an explicit supported subset for incremental corpus processing.

    The historical name is kept for Phase 1--4 compatibility; paths may now be
    PDFs, text-like files, office files, or OCR-supported images.
    """

    if not pdf_paths:
        if not _supported_documents(corpus_root):
            logger.info(
                "corpus_empty",
                extra={
                    "event": "document_loading",
                    "knowledge_root": str(corpus_root),
                    "active_corpus_root": str(corpus_root),
                },
            )
        return []

    effective_config = config or KnowledgeOSConfig(
        project_root=corpus_root.parent.parent
        if corpus_root.name == "files" and corpus_root.parent.name == "data"
        else corpus_root,
        knowledge_root=corpus_root,
    )
    documents: list[Document] = []
    total = len(pdf_paths)
    for position, path in enumerate(pdf_paths, start=1):
        inspection = inspect_ingestion_candidate(path, corpus_root=corpus_root, ocr_engine=effective_config.ocr_engine)
        started = time.perf_counter()
        logger.info(
            "document_indexing_started",
            extra={
                "event": "document_indexing",
                "current_document_number": position,
                "total_documents": total,
                "document_filename": inspection["filename"],
                "relative_path": inspection["relative_path"],
                "extension": inspection["extension"],
                "loader": inspection["loader_selected"],
            },
        )
        try:
            parsed = _load_supported_path(path, effective_config, corpus_root)
        except Exception as exc:
            logger.error(
                "document_indexing_failed",
                extra={
                    "event": "document_indexing",
                    "current_document_number": position,
                    "total_documents": total,
                    "document_filename": inspection["filename"],
                    "relative_path": inspection["relative_path"],
                    "extension": inspection["extension"],
                    "loader": inspection["loader_selected"],
                    "elapsed_seconds": round(time.perf_counter() - started, 3),
                    "error_type": type(exc).__name__,
                },
            )
            raise
        if not parsed:
            logger.warning(
                "document_indexing_skipped",
                extra={
                    "event": "document_indexing",
                    "current_document_number": position,
                    "total_documents": total,
                    "document_filename": inspection["filename"],
                    "relative_path": inspection["relative_path"],
                    "extension": inspection["extension"],
                    "loader": inspection["loader_selected"],
                    "skip_reason": inspection["skip_reason"] or "parser_rejected_file",
                    "elapsed_seconds": round(time.perf_counter() - started, 3),
                },
            )
            continue
        documents.extend(parsed)
        logger.info(
            "document_parsing_completed",
            extra={
                "event": "document_indexing",
                "current_document_number": position,
                "total_documents": total,
                "document_filename": inspection["filename"],
                "relative_path": inspection["relative_path"],
                "extension": inspection["extension"],
                "loader": inspection["loader_selected"],
                "page_count": max((item.metadata.get("page_number") or 0 for item in parsed), default=None) or None,
                "parsed_documents": len(parsed),
                "elapsed_seconds": round(time.perf_counter() - started, 3),
            },
        )
    return documents


def load_all_documents(config: KnowledgeOSConfig) -> list[Document]:
    """Load all supported local document types."""

    return [*load_text_documents(config), *load_pdf_documents(config)]


def summarize_documents(documents: list[Document]) -> dict[str, Any]:
    """Return an inspectable document summary suitable for notebook display."""

    loader_counts = Counter(
        str(document.metadata.get("loader_type", "unknown")) for document in documents
    )
    status_counts = Counter(
        str(document.metadata.get("file_support_status", "unknown"))
        for document in documents
    )
    ocr_counts = Counter(
        str(document.metadata.get("ocr_status"))
        for document in documents
        if document.metadata.get("ocr_status")
    )
    total_characters = sum(len(document.page_content) for document in documents)
    return {
        "document_count": len(documents),
        "total_characters": total_characters,
        "average_characters": round(total_characters / len(documents), 1)
        if documents
        else 0.0,
        "loader_counts": dict(sorted(loader_counts.items())),
        "support_status_counts": dict(sorted(status_counts.items())),
        "ocr_status_counts": dict(sorted(ocr_counts.items())),
        "sources": sorted(
            {str(document.metadata.get("file_name", "")) for document in documents}
        ),
    }


def scan_configured_file_format_readiness(config: KnowledgeOSConfig) -> dict[str, Any]:
    return scan_file_format_readiness(config.knowledge_root)
