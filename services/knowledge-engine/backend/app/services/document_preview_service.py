"""Safe document file, preview, and thumbnail helpers for Corpus documents."""

from __future__ import annotations

import csv
import json
import mimetypes
import re
import uuid
from dataclasses import dataclass
from html import escape, unescape
from io import StringIO
from pathlib import Path
from typing import Any

from fastapi import HTTPException, status
from fastapi.responses import FileResponse

from backend.app.core.config import settings
from backend.app.db.session import SessionLocal
from backend.app.models.knowledge import DocumentChunk
from backend.app.services.document_rendering_service import (
    INLINE_CONVERSION_TARGETS,
    rendered_preview_path,
    viewer_asset_payload,
)
from sqlalchemy import select


TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".html", ".htm", ".json", ".xml", ".yaml", ".yml"}
TABLE_EXTENSIONS = {".csv", ".xlsx", ".xls"}
OFFICE_EXTENSIONS = {".docx", ".doc", ".pptx", ".ppt"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp", ".gif"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | TABLE_EXTENSIONS | OFFICE_EXTENSIONS | IMAGE_EXTENSIONS | {".pdf"}
DIRECT_INLINE_EXTENSIONS = TEXT_EXTENSIONS | {".pdf", ".png", ".jpg", ".jpeg", ".bmp", ".webp", ".gif", ".csv"}
MAX_TEXT_PREVIEW_BYTES = 256 * 1024
MAX_PREVIEW_CHARS = 24_000
THUMBNAIL_SIZE = (360, 240)


@dataclass(frozen=True)
class ResolvedDocument:
    metadata: dict[str, Any]
    path: Path
    extension: str
    content_hash: str


def _safe_name(value: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9_.-]+", "-", value).strip(".-")
    return cleaned or "document"


def _is_within(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
        return True
    except ValueError:
        return False


def _cache_dir(name: str) -> Path:
    directory = settings.outputs_path / name
    directory.mkdir(parents=True, exist_ok=True)
    return directory


def _cache_key(document: ResolvedDocument, page: int | None = None) -> str:
    page_part = f"-p{max(page or 1, 1)}" if page is not None else ""
    return f"{document.metadata['id']}-{_safe_name(document.content_hash)}{page_part}"


def _preview_cache_path(document: ResolvedDocument, page: int | None, chunk_id: str | None) -> Path:
    chunk_part = f"-c{_safe_name(chunk_id)}" if chunk_id else ""
    return _cache_dir("previews") / f"{_cache_key(document, page)}{chunk_part}.json"


def _media_type(path: Path, fallback: str | None = None) -> str:
    guessed, _ = mimetypes.guess_type(path.name)
    return fallback or guessed or "application/octet-stream"


def _read_text_head(path: Path, max_bytes: int = MAX_TEXT_PREVIEW_BYTES) -> str:
    with path.open("rb") as handle:
        data = handle.read(max_bytes + 1)
    suffix = "\n\n[Preview truncated.]" if len(data) > max_bytes else ""
    if len(data) > max_bytes:
        data = data[:max_bytes]
    return data.decode("utf-8", errors="replace")[:MAX_PREVIEW_CHARS] + suffix


def _decode_text(value: str) -> str:
    return unescape(value).replace("\r\n", "\n").replace("\r", "\n")


def _trim_preview_text(value: str) -> str:
    return _decode_text(value)[:MAX_PREVIEW_CHARS]


def _text_lines(value: str, *, limit: int = 6) -> list[str]:
    return [line.strip() for line in _decode_text(value).splitlines() if line.strip()][:limit]


def _html_to_safe_fragment(value: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except Exception:
        return f"<pre>{escape(_trim_preview_text(value))}</pre>"

    soup = BeautifulSoup(value, "html.parser")
    for tag in soup.find_all(["script", "style", "iframe", "object", "embed", "meta", "link"]):
        tag.decompose()
    for tag in soup.find_all(True):
        removable = [
            attribute
            for attribute in list(tag.attrs)
            if attribute.lower().startswith("on")
            or (
                isinstance(tag.attrs.get(attribute), str)
                and str(tag.attrs.get(attribute)).strip().lower().startswith("javascript:")
            )
        ]
        for attribute in removable:
            tag.attrs.pop(attribute, None)
    fragment = soup.body or soup
    return str(fragment)[: MAX_PREVIEW_CHARS * 2]


def _markdown_to_html(value: str) -> str:
    try:
        import markdown
    except Exception:
        return f"<pre>{escape(_trim_preview_text(value))}</pre>"
    rendered = markdown.markdown(
        value,
        extensions=["tables", "fenced_code"],
        output_format="html5",
    )
    return _html_to_safe_fragment(rendered)


def _json_to_pretty(value: str) -> str:
    try:
        return json.dumps(json.loads(value), indent=2, ensure_ascii=False)[:MAX_PREVIEW_CHARS]
    except Exception:
        return _trim_preview_text(value)


def _xml_to_pretty(value: str) -> str:
    try:
        from xml.dom import minidom

        parsed = minidom.parseString(value.encode("utf-8"))
        return parsed.toprettyxml(indent="  ")[:MAX_PREVIEW_CHARS]
    except Exception:
        return _trim_preview_text(value)


def _yaml_to_pretty(value: str) -> str:
    try:
        import yaml

        parsed = yaml.safe_load(value)
        return yaml.safe_dump(parsed, allow_unicode=True, sort_keys=False)[:MAX_PREVIEW_CHARS]
    except Exception:
        return _trim_preview_text(value)


def _csv_rows(path: Path, max_rows: int = 12, max_cols: int = 6) -> list[list[str]]:
    text = _read_text_head(path, 96 * 1024)
    reader = csv.reader(StringIO(text))
    rows: list[list[str]] = []
    for row in reader:
        rows.append([cell.strip()[:80] for cell in row[:max_cols]])
        if len(rows) >= max_rows:
            break
    return rows


def _xlsx_rows(
    path: Path,
    max_rows: int = 12,
    max_cols: int = 6,
) -> tuple[list[list[str]], int, list[str], str | None]:
    try:
        from openpyxl import load_workbook
    except Exception:
        return [], 0, [], None
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        worksheet = workbook.worksheets[0]
        rows: list[list[str]] = []
        for row in worksheet.iter_rows(max_row=max_rows, max_col=max_cols, values_only=True):
            rows.append(["" if cell is None else str(cell)[:80] for cell in row])
        sheet_names = [worksheet.title for worksheet in workbook.worksheets]
        return rows, len(workbook.worksheets), sheet_names, worksheet.title
    finally:
        workbook.close()


def _xls_rows(
    path: Path,
    max_rows: int = 12,
    max_cols: int = 6,
) -> tuple[list[list[str]], int, list[str], str | None]:
    try:
        import xlrd
    except Exception:
        return [], 0, [], None
    try:
        workbook = xlrd.open_workbook(path.as_posix(), on_demand=True)
    except Exception:
        return [], 0, [], None
    try:
        sheet_names = workbook.sheet_names()
        if not sheet_names:
            return [], 0, [], None
        sheet = workbook.sheet_by_index(0)
        rows: list[list[str]] = []
        for row_index in range(min(sheet.nrows, max_rows)):
            row = [
                str(sheet.cell_value(row_index, col_index))[:80].strip()
                for col_index in range(min(sheet.ncols, max_cols))
            ]
            rows.append(row)
        return rows, len(sheet_names), sheet_names, sheet_names[0]
    finally:
        workbook.release_resources()


def _docx_preview(path: Path) -> tuple[str, str, list[list[str]]]:
    try:
        from docx import Document
    except Exception:
        return "", "", []
    try:
        document = Document(path)
    except Exception:
        return "", "", []
    lines: list[str] = []
    html_parts: list[str] = []
    first_table: list[list[str]] = []
    total = 0
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        lines.append(text)
        style_name = str(getattr(paragraph.style, "name", "") or "").casefold()
        if style_name.startswith("heading 1"):
            html_parts.append(f"<h1>{escape(text)}</h1>")
        elif style_name.startswith("heading 2"):
            html_parts.append(f"<h2>{escape(text)}</h2>")
        elif style_name.startswith("heading 3"):
            html_parts.append(f"<h3>{escape(text)}</h3>")
        else:
            html_parts.append(f"<p>{escape(text)}</p>")
        total += len(text)
        if total >= MAX_PREVIEW_CHARS:
            lines.append("[Preview truncated.]")
            html_parts.append("<p><em>Preview truncated.</em></p>")
            break
    for table_index, table in enumerate(document.tables[:3]):
        rows: list[list[str]] = []
        html_parts.append("<table><tbody>")
        for row in table.rows[:8]:
            values = [cell.text.strip()[:120] for cell in row.cells[:6]]
            rows.append(values)
            tag = "th" if len(rows) == 1 else "td"
            html_parts.append(
                "<tr>" + "".join(f"<{tag}>{escape(value)}</{tag}>" for value in values) + "</tr>"
            )
        html_parts.append("</tbody></table>")
        if table_index == 0:
            first_table = rows
    return "\n\n".join(lines)[:MAX_PREVIEW_CHARS], "".join(html_parts), first_table


def _pptx_preview(path: Path) -> tuple[list[dict[str, str]], str]:
    try:
        from pptx import Presentation
    except Exception:
        return [], ""
    try:
        presentation = Presentation(path)
    except Exception:
        return [], ""

    slides: list[dict[str, str]] = []
    preview_lines: list[str] = []
    for index, slide in enumerate(presentation.slides[:12], start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            paragraphs = [
                paragraph.text.strip()
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text and paragraph.text.strip()
            ]
            if not paragraphs:
                continue
            if not title:
                title = paragraphs[0]
                remainder = paragraphs[1:]
            else:
                remainder = paragraphs
            body_parts.extend(remainder[:4])
        body = "\n".join(body_parts[:6]).strip()
        if title or body:
            slides.append(
                {
                    "index": str(index),
                    "title": title or f"Slide {index}",
                    "body": body,
                }
            )
            preview_lines.append(title or f"Slide {index}")
            if body:
                preview_lines.append(body)
    return slides, "\n\n".join(preview_lines)[:MAX_PREVIEW_CHARS]


def _pdf_page_count(path: Path) -> int | None:
    try:
        import fitz
    except Exception:
        return None
    try:
        with fitz.open(path) as document:
            return int(document.page_count)
    except Exception:
        return None


def _pdf_page_text(path: Path, page: int) -> str:
    try:
        import fitz
    except Exception:
        return ""
    try:
        with fitz.open(path) as document:
            index = min(max(page, 1), document.page_count) - 1
            return _trim_preview_text(document.load_page(index).get_text("text"))
    except Exception:
        return ""


def _render_pdf_thumbnail(path: Path, output_path: Path, page: int) -> bool:
    try:
        import fitz
    except Exception:
        return False
    try:
        with fitz.open(path) as document:
            index = min(max(page, 1), document.page_count) - 1
            pixmap = document.load_page(index).get_pixmap(matrix=fitz.Matrix(1.4, 1.4), alpha=False)
            pixmap.save(output_path)
        return True
    except Exception:
        return False


def _draw_card(
    output_path: Path,
    title: str,
    subtitle: str,
    detail: str = "",
    rows: list[list[str]] | None = None,
    lines: list[str] | None = None,
) -> None:
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", THUMBNAIL_SIZE, "#f8faf7")
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()
    strong = ImageFont.load_default()
    draw.rounded_rectangle((1, 1, THUMBNAIL_SIZE[0] - 2, THUMBNAIL_SIZE[1] - 2), radius=16, outline="#dbe5d8", width=2, fill="#ffffff")
    draw.rounded_rectangle((18, 18, 78, 78), radius=12, fill="#edf6e9")
    draw.text((34, 39), (subtitle[:4] or "FILE").upper(), fill="#25611f", font=strong)
    draw.text((96, 23), title[:38], fill="#0f172a", font=strong)
    draw.text((96, 45), subtitle[:42], fill="#64748b", font=font)
    if detail:
        draw.text((96, 64), detail[:42], fill="#2f6d25", font=font)

    if rows:
        x0, y0 = 18, 96
        col_width = 52
        row_height = 18
        for row_index, row in enumerate(rows[:6]):
            for col_index, cell in enumerate(row[:6]):
                x = x0 + col_index * col_width
                y = y0 + row_index * row_height
                fill = "#f8fafc" if row_index == 0 else "#ffffff"
                draw.rectangle((x, y, x + col_width, y + row_height), outline="#e2e8f0", fill=fill)
                draw.text((x + 3, y + 4), str(cell)[:8], fill="#334155", font=font)
    else:
        content_lines = lines or [title[i : i + 42] for i in range(0, min(len(title), 126), 42)]
        for index, line in enumerate(content_lines[:6]):
            draw.text((18, 102 + index * 20), line[:46], fill="#334155", font=font)

    image.save(output_path, format="PNG")


def _render_image_thumbnail(path: Path, output_path: Path) -> bool:
    try:
        from PIL import Image, ImageOps
    except Exception:
        return False
    try:
        with Image.open(path) as image:
            image.seek(0)
            image = ImageOps.exif_transpose(image.convert("RGB"))
            image.thumbnail(THUMBNAIL_SIZE)
            canvas = Image.new("RGB", THUMBNAIL_SIZE, "#f8faf7")
            x = (THUMBNAIL_SIZE[0] - image.width) // 2
            y = (THUMBNAIL_SIZE[1] - image.height) // 2
            canvas.paste(image, (x, y))
            canvas.save(output_path, format="PNG")
        return True
    except Exception:
        return False


def resolve_document(metadata: dict[str, Any]) -> ResolvedDocument:
    relative_path = str(metadata.get("relative_path") or "").replace("\\", "/").strip("/")
    candidate = Path(relative_path)
    if not relative_path or candidate.is_absolute() or ".." in candidate.parts:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid document path metadata.")

    root = settings.data_files_path.resolve()
    path = (root / candidate).resolve()
    if not _is_within(path, root) or not path.is_file():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document file not found.")

    extension = str(metadata.get("extension") or path.suffix).casefold()
    content_hash = str(metadata.get("content_hash") or f"{path.stat().st_mtime_ns}-{path.stat().st_size}")
    return ResolvedDocument(metadata=metadata, path=path, extension=extension, content_hash=content_hash)


def _converted_document(document: ResolvedDocument, output_format: str) -> ResolvedDocument | None:
    converted_path = rendered_preview_path(document, output_format)
    if converted_path is None or not converted_path.is_file():
        return None
    metadata = dict(document.metadata)
    metadata["extension"] = f".{output_format.lstrip('.')}"
    metadata["mime_type"] = "application/pdf" if output_format == "pdf" else metadata.get("mime_type")
    return ResolvedDocument(
        metadata=metadata,
        path=converted_path,
        extension=f".{output_format.lstrip('.')}",
        content_hash=f"{document.content_hash}-{output_format}",
    )


def _chunk_context(document: ResolvedDocument, chunk_id: str | None) -> dict[str, Any]:
    if SessionLocal is None or not chunk_id:
        return {}
    try:
        document_uuid = uuid.UUID(str(document.metadata["id"]))
    except (KeyError, TypeError, ValueError):
        return {}
    with SessionLocal() as session:
        chunk = session.scalar(
            select(DocumentChunk).where(
                DocumentChunk.document_id == document_uuid,
                DocumentChunk.chunk_id == chunk_id,
            )
        )
    if chunk is None:
        return {}
    return {
        "page": chunk.page,
        "highlight_text": _decode_text(str(chunk.text_preview or "")),
    }


def file_response(
    document: ResolvedDocument,
    *,
    disposition: str = "inline",
) -> FileResponse:
    return FileResponse(
        document.path,
        media_type=_media_type(document.path, document.metadata.get("mime_type")),
        filename=document.metadata.get("name") or document.path.name,
        content_disposition_type=disposition,
    )


def view_response(document: ResolvedDocument) -> FileResponse:
    if document.extension in DIRECT_INLINE_EXTENSIONS:
        return file_response(document, disposition="inline")

    target_format = INLINE_CONVERSION_TARGETS.get(document.extension)
    if target_format:
        converted = _converted_document(document, target_format)
        if converted is not None:
            return file_response(converted, disposition="inline")
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Inline preview is not available for this document. Use the workspace preview or Download to inspect the original file.",
        )

    raise HTTPException(
        status_code=status.HTTP_404_NOT_FOUND,
        detail="Inline preview is not available for this document type. Use the workspace preview or Download to inspect the original file.",
    )


def thumbnail_response(document: ResolvedDocument, page: int | None = None) -> FileResponse:
    page_number = max(page or 1, 1)
    output_path = _cache_dir("thumbnails") / f"{_cache_key(document, page_number)}.png"
    if not output_path.is_file():
        rendered = False
        thumbnail_source = document
        if document.extension in INLINE_CONVERSION_TARGETS:
            converted = _converted_document(document, INLINE_CONVERSION_TARGETS[document.extension])
            if converted is not None:
                thumbnail_source = converted

        if thumbnail_source.extension == ".pdf":
            rendered = _render_pdf_thumbnail(thumbnail_source.path, output_path, page_number)
        elif document.extension in IMAGE_EXTENSIONS:
            rendered = _render_image_thumbnail(document.path, output_path)

        if not rendered:
            rows: list[list[str]] | None = None
            lines: list[str] | None = None
            if thumbnail_source.extension == ".csv":
                rows = _csv_rows(document.path)
            elif thumbnail_source.extension == ".xlsx":
                rows, _, _, _ = _xlsx_rows(document.path)
            elif thumbnail_source.extension == ".xls":
                rows, _, _, _ = _xls_rows(document.path)
            else:
                preview = _preview_content(thumbnail_source, page=page_number)
                if preview.get("slides"):
                    first_slide = preview["slides"][0]
                    lines = _text_lines(
                        f"{first_slide.get('title', '')}\n{first_slide.get('body', '')}",
                        limit=6,
                    )
                else:
                    lines = _text_lines(str(preview.get("preview_text") or ""), limit=6)
            try:
                _draw_card(
                    output_path,
                    str(document.metadata.get("name") or document.path.name),
                    document.extension.replace(".", "") or "document",
                    str(document.metadata.get("indexing_status") or "pending"),
                    rows=rows,
                    lines=lines,
                )
            except Exception:
                raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Thumbnail generation failed.")
    return FileResponse(output_path, media_type="image/png")


def _preview_content(document: ResolvedDocument, *, page: int | None = None) -> dict[str, Any]:
    preview_text = ""
    rows: list[list[str]] = []
    slides: list[dict[str, str]] = []
    extraction_method = "metadata"
    render_kind = "card"
    extra: dict[str, Any] = {
        "rendered_html": None,
        "sheet_names": [],
        "active_sheet": None,
        "slides": slides,
    }

    if document.extension in TEXT_EXTENSIONS:
        preview_text = _trim_preview_text(_read_text_head(document.path))
        extraction_method = "text_head"
        if document.extension == ".json":
            preview_text = _json_to_pretty(preview_text)
            render_kind = "code"
        elif document.extension == ".xml":
            preview_text = _xml_to_pretty(preview_text)
            render_kind = "code"
        elif document.extension in {".yaml", ".yml"}:
            preview_text = _yaml_to_pretty(preview_text)
            render_kind = "code"
        elif document.extension in {".md", ".markdown"}:
            render_kind = "markdown"
            extra["rendered_html"] = _markdown_to_html(preview_text)
        elif document.extension in {".html", ".htm"}:
            render_kind = "html"
            extra["rendered_html"] = _html_to_safe_fragment(preview_text)
        else:
            render_kind = "text"
    elif document.extension == ".csv":
        rows = _csv_rows(document.path)
        preview_text = _trim_preview_text("\n".join([", ".join(row) for row in rows]))
        extraction_method = "csv_head"
        render_kind = "table"
    elif document.extension == ".xlsx":
        rows, sheet_count, sheet_names, active_sheet = _xlsx_rows(document.path)
        preview_text = _trim_preview_text("\n".join(["\t".join(row) for row in rows]))
        extraction_method = "spreadsheet_head" if rows else "metadata"
        render_kind = "spreadsheet" if rows else "card"
        extra["sheet_count"] = sheet_count
        extra["sheet_names"] = sheet_names
        extra["active_sheet"] = active_sheet
    elif document.extension == ".xls":
        rows, sheet_count, sheet_names, active_sheet = _xls_rows(document.path)
        preview_text = _trim_preview_text("\n".join(["\t".join(row) for row in rows]))
        extraction_method = "spreadsheet_head" if rows else "metadata"
        render_kind = "spreadsheet" if rows else "card"
        extra["sheet_count"] = sheet_count
        extra["sheet_names"] = sheet_names
        extra["active_sheet"] = active_sheet
    elif document.extension == ".docx":
        preview_text, rendered_html, first_table = _docx_preview(document.path)
        extraction_method = "docx_text" if preview_text else "metadata"
        render_kind = "docx" if preview_text or rendered_html else "card"
        extra["rendered_html"] = rendered_html or None
        if first_table:
            rows = first_table
    elif document.extension == ".pdf":
        render_kind = "pdf"
        extraction_method = "pdf_page_text"
        extra["page_count"] = document.metadata.get("page_count") or _pdf_page_count(document.path)
        preview_text = _pdf_page_text(document.path, page or 1)
    elif document.extension in IMAGE_EXTENSIONS:
        render_kind = "image"
        extraction_method = "file_stream"
    elif document.extension == ".pptx":
        slides, preview_text = _pptx_preview(document.path)
        extraction_method = "pptx_slides" if slides else "metadata"
        render_kind = "slides" if slides else "card"
        extra["slides"] = slides
    elif document.extension in OFFICE_EXTENSIONS:
        render_kind = "card"

    return {
        "preview_text": preview_text,
        "table_rows": rows,
        "render_kind": render_kind,
        "extraction_method": extraction_method,
        **extra,
    }


def preview_payload(document: ResolvedDocument, *, page: int | None = None, chunk_id: str | None = None) -> dict[str, Any]:
    chunk_context = _chunk_context(document, chunk_id)
    viewer_asset = viewer_asset_payload(document)
    preview_source = document
    if (
        document.extension in INLINE_CONVERSION_TARGETS
        and viewer_asset.get("viewer_ready")
        and viewer_asset.get("viewer_format") == "pdf"
    ):
        converted = _converted_document(document, "pdf")
        if converted is not None:
            preview_source = converted

    resolved_page = page or chunk_context.get("page")
    if resolved_page is not None:
        resolved_page = max(int(resolved_page), 1)
    cache_path = _preview_cache_path(preview_source, resolved_page, chunk_id)
    cached: dict[str, Any] | None = None
    if cache_path.is_file():
        try:
            cached = json.loads(cache_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            cached = None
    if cached is None:
        cached = _preview_content(preview_source, page=resolved_page)
        try:
            cache_path.write_text(json.dumps(cached, ensure_ascii=False), encoding="utf-8")
        except OSError:
            pass

    preview_text = _trim_preview_text(str(cached.get("preview_text") or ""))
    rows = cached.get("table_rows") if isinstance(cached.get("table_rows"), list) else []
    highlight_text = _decode_text(str(chunk_context.get("highlight_text") or "")).strip()
    if not highlight_text and preview_text:
        highlight_text = preview_text[:1000]
    if not highlight_text:
        highlight_text = f"{document.metadata.get('name') or document.path.name} is available as a corpus file. Inline extracted text is not available for this format."

    file_id = str(document.metadata["id"])
    return {
        **document.metadata,
        "preview_text": preview_text,
        "highlight_text": _decode_text(highlight_text),
        "page": resolved_page,
        "chunk_id": chunk_id,
        "document_id": file_id,
        "relative_path": document.metadata.get("relative_path"),
        "page_count": cached.get("page_count") or document.metadata.get("page_count"),
        "open_url": f"/api/corpus/document/{file_id}/view",
        "download_url": f"/api/corpus/document/{file_id}/download",
        "file_url": f"/api/corpus/document/{file_id}/file",
        "thumbnail_url": f"/api/corpus/document/{file_id}/thumbnail?page={max(resolved_page or 1, 1)}",
        "read_error": None,
        "render_kind": cached.get("render_kind") or "card",
        "extraction_method": cached.get("extraction_method") or "metadata",
        "table_rows": rows,
        "supported_preview": document.extension in SUPPORTED_EXTENSIONS or bool(viewer_asset.get("viewer_ready")),
        **viewer_asset,
        **{key: value for key, value in cached.items() if key not in {"preview_text", "table_rows", "render_kind", "extraction_method"}},
    }


def parse_document_id(document_id: str) -> uuid.UUID:
    try:
        return uuid.UUID(document_id)
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid document id.") from exc
